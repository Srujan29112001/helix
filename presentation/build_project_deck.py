"""Builds the Helix capstone *project deck* — 22 slides, dark brand-matched, for a
mixed audience (non-tech -> learners -> data scientists -> professors -> business
leaders).

Design takes the user's "HELIX — The Autonomous Data Scientist" deck as inspiration
(full-bleed hero art on title/section slides, before/after, deep dives, self-healing
loop, prompt-engineering, security, resilience, deployment comparison) and merges it
with accurate, current technical content + REAL app screenshots.

Images auto-embed from presentation/assets/ by filename; a labelled placeholder is
drawn for anything missing. See presentation/IMAGE_GUIDE.md.

Run:  python presentation/build_project_deck.py  ->  presentation/Helix_Project_Deck.pptx
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

try:
    from PIL import Image
except Exception:
    Image = None

# ----------------------------- palette -----------------------------
INK = RGBColor(0x04, 0x06, 0x0F)
INK2 = RGBColor(0x09, 0x0E, 0x1A)
PANEL = RGBColor(0x0C, 0x12, 0x24)
PANEL2 = RGBColor(0x10, 0x18, 0x2E)
EDGE = RGBColor(0x1A, 0x24, 0x40)
CYAN = RGBColor(0x25, 0xD7, 0xF0)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
AZURE = RGBColor(0x38, 0xBD, 0xF8)
IRIS = RGBColor(0xA7, 0x8B, 0xFA)
MAGENTA = RGBColor(0xE8, 0x79, 0xF9)
ACID = RGBColor(0x9A, 0xE6, 0x4A)
GOLD = RGBColor(0xFB, 0xBF, 0x24)
CORAL = RGBColor(0xFB, 0x71, 0x85)
TEAL = RGBColor(0x2D, 0xD4, 0xBF)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MIST = RGBColor(0xC9, 0xD6, 0xEF)
MUTE = RGBColor(0x76, 0x85, 0x9F)

HEAD = "Segoe UI Semibold"
BODY = "Segoe UI"
MONO = "Consolas"
CENTER, LEFT, RIGHT = PP_ALIGN.CENTER, PP_ALIGN.LEFT, PP_ALIGN.RIGHT
MID, TOP = MSO_ANCHOR.MIDDLE, MSO_ANCHOR.TOP

HERE = os.path.dirname(__file__)
ASSETS = os.path.join(HERE, "assets")
TOTAL = 22

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5

AGENTS9 = [
    ("Planner", VIOLET), ("Coder", CYAN), ("Executor", AZURE), ("Critic", GOLD),
    ("AutoML", ACID), ("Explainer", MAGENTA), ("Visualizer", TEAL),
    ("Researcher", IRIS), ("Reporter", CORAL),
]


# ----------------------------- helpers -----------------------------
def slide(bg=INK):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def text(s, x, y, w, h, runs, align=LEFT, anchor=TOP, space=2):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        for t, size, color, bold, font in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
    return tb


def one(s, x, y, w, h, t, size, color, bold=False, font=BODY, align=LEFT, anchor=TOP):
    return text(s, x, y, w, h, [[(t, size, color, bold, font)]], align=align, anchor=anchor)


def box(s, x, y, w, h, fill=PANEL, line=EDGE, line_w=1.0, radius=0.08):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def frame(s, x, y, w, h, color, w_pt=1.25):
    shp = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.fill.background()
    shp.line.color.rgb = color
    shp.line.width = Pt(w_pt)
    shp.shadow.inherit = False
    return shp


def dot(s, x, y, d, color):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    c.shadow.inherit = False
    return c


def eyebrow(s, x, y, label, color=CYAN):
    dot(s, x, y + 0.05, 0.12, color)
    one(s, x + 0.22, y - 0.05, 9, 0.4, label.upper(), 12, MUTE, True, MONO)


def title(s, x, y, t, w=11.6, size=32, color=WHITE):
    one(s, x, y, w, 1.1, t, size, color, True, HEAD)


def footer(s, n):
    one(s, 1.0, 7.06, 8, 0.3, "HELIX · The Autonomous Data Scientist", 9, MUTE, False, MONO)
    one(s, 10.3, 7.06, 2.03, 0.3, f"{n:02d} / {TOTAL}", 9, MUTE, False, MONO, align=RIGHT)


def _ratio(fname, default=1.6):
    """Native width/height ratio of an asset (so images never get distorted)."""
    if Image is None:
        return default
    try:
        with Image.open(os.path.join(ASSETS, fname)) as im:
            return im.size[0] / im.size[1]
    except Exception:
        return default


def place_w(s, x, y, w, fname, tag="IMAGE", caption="", accent=CYAN, border=True):
    """Landscape image: fix the width, derive the height from the native ratio."""
    full = os.path.join(ASSETS, fname)
    if os.path.exists(full):
        h = w / _ratio(fname, 1.6)
        s.shapes.add_picture(full, Inches(x), Inches(y), Inches(w), Inches(h))
        if border:
            frame(s, x, y, w, h, EDGE, 1.0)
        return h
    h = w / 1.6
    box(s, x, y, w, h, INK2, accent, 1.25)
    one(s, x + 0.2, y + 0.18, w - 0.4, 0.3, fname, 9.5, accent, True, MONO)
    one(s, x, y + h / 2 - 0.3, w, 0.4, tag, 13, accent, True, MONO, CENTER, MID)
    one(s, x + 0.4, y + h / 2 + 0.15, w - 0.8, 0.5, caption, 11, MUTE, False, BODY, CENTER, MID)
    return h


def place_box(s, x, y, w, h, fname, tag="IMAGE", caption="", accent=CYAN):
    """Fit an image inside (x,y,w,h) preserving ratio, centered (letterboxed)."""
    full = os.path.join(ASSETS, fname)
    if not os.path.exists(full):
        box(s, x, y, w, h, INK2, accent, 1.25)
        one(s, x, y + h / 2 - 0.3, w, 0.4, tag, 13, accent, True, MONO, CENTER, MID)
        one(s, x + 0.4, y + h / 2 + 0.15, w - 0.8, 0.5, caption, 11, MUTE, False, BODY, CENTER, MID)
        return
    r = _ratio(fname, 1.6)
    iw, ih = w, w / r
    if ih > h:
        ih, iw = h, h * r
    ix, iy = x + (w - iw) / 2, y + (h - ih) / 2
    s.shapes.add_picture(full, Inches(ix), Inches(iy), Inches(iw), Inches(ih))
    frame(s, ix, iy, iw, ih, EDGE, 1.0)


def hero_right(s, fname, tag="HERO ART", accent=VIOLET):
    """Portrait image bleeding off the right edge at native ratio; text on the left."""
    r = _ratio(fname, 0.667)
    width = SH * r  # full height, ratio-correct width
    x = SW - width
    if os.path.exists(os.path.join(ASSETS, fname)):
        s.shapes.add_picture(os.path.join(ASSETS, fname), Inches(x), Inches(0), Inches(width), Inches(SH))
    else:
        box(s, x, 0, width, SH, INK2, accent, 1.25, radius=0.0)
        one(s, x, 3.2, width, 0.4, tag, 13, accent, True, MONO, CENTER, MID)
        one(s, x + 0.4, 3.7, width - 0.8, 0.5, fname, 10.5, MUTE, False, MONO, CENTER, MID)
    box(s, x - 0.02, 0, 0.04, SH, EDGE, None)
    return x  # left edge of the image — keep text left of this


def bigstat(s, x, y, w, v, l, c):
    box(s, x, y, w, 1.5, PANEL, EDGE)
    one(s, x, y + 0.16, w, 0.8, v, 34, c, True, HEAD, CENTER)
    one(s, x, y + 1.0, w, 0.4, l, 11.5, MUTE, False, MONO, CENTER)


def checkrow(s, x, y, w, label, helix=True, other="✕", oc=CORAL):
    one(s, x, y, w - 3.0, 0.4, label, 12.5, MIST, False, BODY, anchor=MID)
    one(s, x + w - 2.7, y, 1.2, 0.4, "✓", 14, ACID, True, BODY, CENTER, MID)
    one(s, x + w - 1.4, y, 1.3, 0.4, other, 12.5, oc, other == "✓", BODY, CENTER, MID)


# ============================ 1 — Title ============================
s = slide(INK)
hero_right(s, "art_hero.jpg", tag="NANOBANANA — hero", accent=CYAN)
for i, (_, c) in enumerate(AGENTS9):
    dot(s, 1.0 + i * 0.40, 1.55, 0.2, c)
one(s, 1.02, 2.15, 7.0, 0.4, "CAPSTONE PROJECT  ·  IIIT-HYDERABAD  ·  JUNE 2026", 12, MUTE, True, MONO)
text(s, 1.0, 2.55, 7.0, 1.4, [[("Helix", 76, WHITE, True, HEAD), (".", 76, CYAN, True, HEAD)]])
one(s, 1.02, 3.95, 6.8, 0.5, "The Autonomous Data Scientist", 24, CYAN, True, BODY)
one(s, 1.04, 4.7, 6.6, 1.2, "A nine-agent AI system that takes a CSV and a plain-English goal, then plans, codes, executes, self-heals and reports — turning raw data into a board-ready report in minutes, not days.", 14, MIST, False, BODY)
chips = [("9 AI AGENTS", CYAN), ("SELF-HEALING CODE", GOLD), ("100% EXPLAINABLE", ACID)]
cx = 1.0
for t, c in chips:
    w = 0.30 + len(t) * 0.105
    box(s, cx, 6.05, w, 0.5, INK2, c)
    one(s, cx, 6.05, w, 0.5, t, 11, c, True, MONO, CENTER, MID)
    cx += w + 0.25
one(s, 1.02, 6.85, 7.0, 0.4, "helix-henna.vercel.app    ·    github.com/Srujan29112001/helix", 11.5, MUTE, False, MONO)

# ============================ 2 — Overview ============================
s = slide(INK)
eyebrow(s, 1.0, 0.8, "Overview")
title(s, 1.0, 1.15, "What is Helix?")
one(s, 1.0, 2.0, 7.1, 1.3, "An autonomous, multi-agent data-science studio. Describe your goal in plain English, point it at a CSV, and nine specialized AI agents collaborate to clean, analyze, model, visualize and explain your data — with zero manual coding.", 14.5, MIST, False, BODY)
steps = [("You provide", "A CSV + a goal in the Studio", CYAN),
         ("AI processes", "9 agents analyze & interpret", VIOLET),
         ("You receive", "Charts, report & strategy", ACID)]
for i, (h, b, c) in enumerate(steps):
    y = 3.45 + i * 1.08
    box(s, 1.0, y, 7.1, 0.95, PANEL, EDGE)
    dot(s, 1.3, y + 0.36, 0.24, c)
    one(s, 1.75, y + 0.14, 3.0, 0.4, h, 14, WHITE, True, HEAD)
    one(s, 1.75, y + 0.52, 6.0, 0.4, b, 12, MUTE, False, BODY)
    if i < 2:
        one(s, 4.4, y + 0.92, 0.4, 0.3, "↓", 15, MUTE, True, BODY, CENTER)
place_w(s, 8.35, 2.0, 4.48, "05_landing.png", "SCREENSHOT", "The live Studio (Vercel)", CYAN)
one(s, 8.35, 4.72, 4.48, 0.3, "Four sample datasets + upload — the live product.", 10, MUTE, False, BODY)
stats = [("9", "AI agents", CYAN), ("8", "LLM providers", VIOLET), ("16+", "SVG charts", ACID), ("Dual", "sandbox", GOLD)]
for i, (v, l, c) in enumerate(stats):
    x = 8.35 + (i % 2) * 2.32
    y = 5.2 + (i // 2) * 0.92
    box(s, x, y, 2.16, 0.8, PANEL, EDGE)
    one(s, x + 0.18, y + 0.1, 2.0, 0.45, v, 20, c, True, HEAD)
    one(s, x + 0.18, y + 0.5, 2.0, 0.25, l, 9.5, MUTE, False, MONO)
footer(s, 2)

# ============================ 3 — Problem ============================
s = slide(INK)
hero_right(s, "art_problem.jpg", tag="HERO — problem", accent=CORAL)
eyebrow(s, 1.0, 0.8, "The problem", CORAL)
title(s, 1.0, 1.15, "Why Helix exists", w=6.9)
text(s, 1.0, 2.0, 7.0, 0.9, [[("60–70%", 30, GOLD, True, HEAD),
     ("  of a data scientist's time is lost to repetitive prep — loading, cleaning, encoding — not analysis.", 14, MIST, False, BODY)]])
probs = [
    ("The bottleneck", "Business teams queue for days; insights arrive too late to act on.", CORAL),
    ("The repetitive grind", "Every dataset demands the same boilerplate pipeline, rewritten by hand.", GOLD),
    ("Brittle iteration", "When goals shift or models fail, the whole pipeline is rebuilt manually.", IRIS),
]
for i, (h, b, c) in enumerate(probs):
    y = 3.3 + i * 1.18
    box(s, 1.0, y, 7.0, 1.05, PANEL, EDGE)
    dot(s, 1.3, y + 0.4, 0.22, c)
    one(s, 1.72, y + 0.15, 6.1, 0.4, h, 15, WHITE, True, HEAD)
    one(s, 1.72, y + 0.55, 6.1, 0.4, b, 12, MUTE, False, BODY)
footer(s, 3)

# ============================ 4 — Solution (before/after) ============================
s = slide(INK)
eyebrow(s, 1.0, 0.8, "The solution", ACID)
title(s, 1.0, 1.15, "Days of manual work → minutes of autonomy")
before = ["Manual pipeline takes days to weeks", "Data-scientist time is the bottleneck",
          "Non-technical stakeholders locked out", "Every iteration means manual rework",
          "Results arrive after the context shifts"]
after = ["Autonomous pipeline finishes in minutes", "Plain-English goal → board-ready report",
         "Anyone can run a full ML analysis", "Re-run instantly with a one-line change",
         "Self-healing agents fix their own errors"]
box(s, 1.0, 2.2, 5.55, 4.0, PANEL, CORAL, 1.25)
one(s, 1.3, 2.42, 5.0, 0.4, "BEFORE HELIX", 13, CORAL, True, MONO)
for i, t in enumerate(before):
    y = 3.05 + i * 0.62
    one(s, 1.3, y, 0.4, 0.4, "✕", 13, CORAL, True, BODY, anchor=MID)
    one(s, 1.7, y, 4.7, 0.5, t, 12.5, MIST, False, BODY, anchor=MID)
box(s, 6.78, 2.2, 5.55, 4.0, PANEL, ACID, 1.25)
one(s, 7.08, 2.42, 5.0, 0.4, "AFTER HELIX", 13, ACID, True, MONO)
for i, t in enumerate(after):
    y = 3.05 + i * 0.62
    one(s, 7.08, y, 0.4, 0.4, "✓", 13, ACID, True, BODY, anchor=MID)
    one(s, 7.48, y, 4.7, 0.5, t, 12.5, MIST, False, BODY, anchor=MID)
one(s, 1.0, 6.45, 11.3, 0.4, "Helix replaces days of manual work with minutes of autonomous, explainable analysis.", 13, MUTE, False, BODY, CENTER)
footer(s, 4)

# ============================ 5 — Architecture (diagram) ============================
s = slide(INK)
eyebrow(s, 1.0, 0.72, "System architecture", VIOLET)
title(s, 1.0, 1.06, "A three-tier, stateless pipeline")
place_box(s, 1.0, 1.9, 11.33, 4.6, "diagram_architecture.png", "DIAGRAM", "Frontend · Backend API (9 agents + LangGraph) · AI/ML layer", VIOLET)
one(s, 1.0, 6.65, 11.3, 0.4, "Stateless · no database · in-memory per request · agents stream live to the UI via Server-Sent Events.", 12.5, MUTE, False, BODY, CENTER)
footer(s, 5)

# ============================ 6 — Nine-agent pipeline ============================
s = slide(INK)
eyebrow(s, 1.0, 0.8, "Core architecture", CYAN)
title(s, 1.0, 1.15, "The nine-agent pipeline")
jobs = ["Numbered analysis plan from your goal", "Writes Python via a RAG-grounded LLM",
        "Runs the code in a secure sandbox", "Debugs & fixes errors (≤5 retries)",
        "FLAML searches for the best ML model", "SHAP feature importance + direction",
        "Picks the best of 16+ chart types", "Live web search for domain context",
        "Writes the business report"]
for i, ((name, c), job) in enumerate(zip(AGENTS9, jobs)):
    col, row = i % 3, i // 3
    x = 1.0 + col * 3.83
    y = 2.4 + row * 1.42
    box(s, x, y, 3.55, 1.25, PANEL, EDGE)
    one(s, x + 0.28, y + 0.18, 0.5, 0.4, str(i + 1), 16, c, True, HEAD)
    one(s, x + 0.78, y + 0.2, 2.6, 0.4, name, 14, WHITE, True, HEAD)
    one(s, x + 0.3, y + 0.66, 3.1, 0.5, job, 11, MUTE, False, BODY)
one(s, 1.0, 6.75, 11.3, 0.4, "FLOW   Planner → Coder → Executor → Critic → AutoML → Explainer → Visualizer → Researcher → Reporter", 11.5, CYAN, True, MONO, CENTER)
footer(s, 6)

# ============================ 7 — Deep dive part 1 ============================
s = slide(INK2)
eyebrow(s, 1.0, 0.8, "Deep dive · part 1", AZURE)
title(s, 1.0, 1.15, "Plan · Code · Execute · Critique")
dd = [
    ("1  Planner", "Reads your goal and data schema, then designs a step-by-step analysis plan. Planning first prevents wasted computation downstream.", VIOLET),
    ("2  Coder", "Converts the plan into executable Python (pandas, scikit-learn) — grounded by RAG over library docs to cut hallucinated APIs.", CYAN),
    ("3  Executor", "Runs the generated code inside an isolated sandbox (RestrictedPython / E2B). Untrusted code never touches the host.", AZURE),
    ("4  Critic", "Audits outputs for errors. On failure it feeds a correction back to the Coder — the self-healing loop, up to 5 retries.", GOLD),
]
for i, (h, b, c) in enumerate(dd):
    x = 1.0 + (i % 2) * 5.85
    y = 2.3 + (i // 2) * 2.05
    box(s, x, y, 5.5, 1.8, PANEL, EDGE)
    box(s, x, y, 0.12, 1.8, c, None, radius=0.0)
    one(s, x + 0.4, y + 0.24, 5.0, 0.4, h, 16, c, True, HEAD)
    one(s, x + 0.4, y + 0.78, 4.95, 0.95, b, 12.5, MUTE, False, BODY)
footer(s, 7)

# ============================ 8 — Deep dive part 2 ============================
s = slide(INK2)
eyebrow(s, 1.0, 0.8, "Deep dive · part 2", MAGENTA)
title(s, 1.0, 1.18, "Model · Explain · Visualize · Research · Report", size=25)
dd2 = [
    ("5  AutoML", "FLAML searches LightGBM, XGBoost, RF & Logistic Regression — selecting the best model under a time budget, no manual tuning.", ACID),
    ("6  Explainer", "SHAP turns metrics and feature importances into plain-English drivers any stakeholder understands — with direction.", MAGENTA),
    ("7  Visualizer", "Chooses the best chart from 16+ pure-SVG types; the engine fills every number from real data — a chart can't fabricate a value.", TEAL),
    ("8  Researcher", "RAG + live web search (Tavily / DuckDuckGo) enrich findings with industry benchmarks and domain context.", IRIS),
    ("9  Reporter", "Assembles the deliverable: a 6–7 paragraph executive report with drivers, risks and prioritized recommendations.", CORAL),
]
for i, (h, b, c) in enumerate(dd2):
    if i < 4:
        x = 1.0 + (i % 2) * 5.85
        y = 2.3 + (i // 2) * 1.55
        w = 5.5
    else:
        x, y, w = 1.0, 5.4, 11.35
    box(s, x, y, w, 1.4 if i < 4 else 1.05, PANEL, EDGE)
    box(s, x, y, 0.12, 1.4 if i < 4 else 1.05, c, None, radius=0.0)
    one(s, x + 0.4, y + 0.16, w - 0.7, 0.4, h, 15, c, True, HEAD)
    one(s, x + 0.4, y + 0.6, w - 0.7, 0.7, b, 12, MUTE, False, BODY)
footer(s, 8)

# ============================ 9 — Self-healing loop ============================
s = slide(INK)
eyebrow(s, 1.0, 0.8, "Key innovation", GOLD)
title(s, 1.0, 1.15, "The self-healing loop")
one(s, 1.0, 2.0, 11.2, 0.85, "When generated code fails, Helix doesn't stop. The Critic diagnoses the error and feeds a correction back to the Coder — the pipeline repairs itself, autonomously, up to five times.", 14, MIST, False, BODY)
loop = [("Coder", "writes Python", CYAN), ("Executor", "runs in sandbox", AZURE), ("Critic", "audits output", GOLD)]
for i, (h, b, c) in enumerate(loop):
    x = 1.4 + i * 3.2
    box(s, x, 3.2, 2.6, 1.15, PANEL, c, 1.4)
    one(s, x, 3.36, 2.6, 0.4, h, 16, WHITE, True, HEAD, CENTER)
    one(s, x, 3.82, 2.6, 0.4, b, 11.5, MUTE, False, MONO, CENTER)
    if i < 2:
        one(s, x + 2.55, 3.2, 0.7, 1.15, "→", 22, MUTE, True, BODY, CENTER, MID)
one(s, 7.9, 3.2, 4.3, 1.15, "Pass → AutoML", 14, ACID, True, MONO, anchor=MID)
box(s, 1.4, 4.55, 9.0, 0.6, INK2, GOLD)
one(s, 1.4, 4.55, 9.0, 0.6, "✕ error  →  diagnosis  →  corrected code  →  retry   (max 5×)", 13, GOLD, True, MONO, CENTER, MID)
hl = [("≤ 5", "autonomous retries per step", CYAN), ("0", "manual debugging required", ACID), ("100%", "responses always complete", GOLD)]
for i, (v, l, c) in enumerate(hl):
    x = 1.0 + i * 3.83
    box(s, x, 5.45, 3.55, 1.2, PANEL, EDGE)
    one(s, x, 5.56, 3.55, 0.6, v, 28, c, True, HEAD, CENTER)
    one(s, x, 6.25, 3.55, 0.35, l, 11, MUTE, False, MONO, CENTER)
footer(s, 9)

# ============================ 10 — The Studio (screenshot) ============================
s = slide(INK)
eyebrow(s, 1.0, 0.8, "Technology · frontend", CYAN)
title(s, 1.0, 1.15, "The Studio — set up in four fields")
place_w(s, 1.0, 2.45, 7.1, "06_studio.png", "SCREENSHOT", "Studio setup — sample/upload, goal, target, task, AI engine", CYAN)
fr = [("Dataset", "Sample, or upload CSV / Excel / Parquet / JSON, or paste a URL / Sheet."),
      ("Goal", "Your question in plain English — feeds the Planner & Reporter."),
      ("Target + Task", "Column to predict (or Auto) and task (or Auto) — hides for clustering."),
      ("AI engine", "Provider, model, per-agent temperature & key — keys stay in your browser.")]
for i, (h, b) in enumerate(fr):
    y = 2.15 + i * 1.2
    box(s, 8.4, y, 3.93, 1.07, PANEL, EDGE)
    one(s, 8.62, y + 0.13, 3.6, 0.4, h, 13, CYAN, True, HEAD)
    one(s, 8.62, y + 0.5, 3.6, 0.5, b, 10, MUTE, False, BODY)
footer(s, 10)

# ============================ 11 — Live pipeline (screenshot) ============================
s = slide(INK)
eyebrow(s, 1.0, 0.8, "Watch it work", AZURE)
title(s, 1.0, 1.15, "Full transparency — every step, live")
place_w(s, 1.0, 2.45, 7.1, "10_pipeline.png", "SCREENSHOT", "Pipeline — agents stream code, stdout & the self-heal loop", AZURE)
pp = [("Real-time SSE", "An asyncio.Queue event bus pushes each agent's status the instant it happens — sub-second latency."),
      ("Colab-style", "Complete code, stdout and tracebacks — never truncated."),
      ("Prompt & logic", "Each agent's exact system prompt + the sandbox rules, on tap."),
      ("See the heal", "The Critic's diagnosis and corrected code, in real time.")]
for i, (h, b) in enumerate(pp):
    y = 2.15 + i * 1.2
    box(s, 8.4, y, 3.93, 1.07, PANEL, EDGE)
    one(s, 8.62, y + 0.12, 3.6, 0.4, h, 13, AZURE, True, HEAD)
    one(s, 8.62, y + 0.49, 3.6, 0.5, b, 10, MUTE, False, BODY)
footer(s, 11)

# ============================ 12 — Results (screenshot) ============================
s = slide(INK)
eyebrow(s, 1.0, 0.8, "The payoff", MAGENTA)
title(s, 1.0, 1.15, "A results dashboard anyone can read")
place_w(s, 1.0, 2.45, 7.1, "11_results_telco.png", "SCREENSHOT", "Results on Telco Churn — verdict, metrics, key findings", MAGENTA)
rr = [("Honest verdict", "Grades the signal from strong to near-chance — weak data labels itself."),
      ("Interactive charts", "Hover any point or line for the exact value; real numeric axes."),
      ("Key findings", "Best model, top driver, key segment, correlations — at a glance."),
      ("Business report", "Narrative + recommendation, grounded in live web context.")]
for i, (h, b) in enumerate(rr):
    y = 2.15 + i * 1.2
    box(s, 8.4, y, 3.93, 1.07, PANEL, EDGE)
    one(s, 8.62, y + 0.12, 3.6, 0.4, h, 13, MAGENTA, True, HEAD)
    one(s, 8.62, y + 0.49, 3.6, 0.5, b, 10, MUTE, False, BODY)
footer(s, 12)

# ============================ 13 — Backend ============================
s = slide(INK)
eyebrow(s, 1.0, 0.8, "Technology · backend", AZURE)
title(s, 1.0, 1.15, "Where the intelligence lives")
groups = [
    ("FastAPI + Uvicorn", "Async Python, native SSE streaming.", AZURE),
    ("LangGraph", "9-agent StateGraph + self-heal loops.", VIOLET),
    ("FLAML AutoML", "Fast model search under a budget.", ACID),
    ("SHAP", "Exact, game-theoretic explainability.", MAGENTA),
    ("ChromaDB + MiniLM", "In-memory RAG over DS recipes.", CYAN),
    ("RestrictedPython + E2B", "Dual sandbox — jail or microVM.", CORAL),
]
for i, (h, b, c) in enumerate(groups):
    x = 1.0 + (i % 2) * 3.25
    y = 2.3 + (i // 2) * 1.5
    box(s, x, y, 3.05, 1.32, PANEL, EDGE)
    dot(s, x + 0.26, y + 0.34, 0.18, c)
    one(s, x + 0.58, y + 0.19, 2.4, 0.4, h, 11.5, WHITE, True, HEAD)
    one(s, x + 0.28, y + 0.68, 2.7, 0.55, b, 10, MUTE, False, BODY)
place_box(s, 7.65, 2.3, 4.95, 3.55, "art_backend.jpg", "HERO — backend", "the engine", AZURE)
one(s, 7.65, 6.0, 4.95, 0.75, "Stateless, async Python — the core analysis engine powers cleaning, modeling, explainability and the chart gallery.", 11, MUTE, False, BODY)
footer(s, 13)

# ============================ 14 — Data flow ============================
s = slide(INK2)
eyebrow(s, 1.0, 0.8, "How it works", TEAL)
title(s, 1.0, 1.15, "End-to-end data flow")
flow = [
    ("1  User input", "Upload a CSV, set a plain-English goal, configure the LLM per agent role.", CYAN),
    ("2  Data cleaning", "Strip, coerce numerics, impute missing values, label-encode, parse dates.", AZURE),
    ("3  Agent pipeline", "LangGraph runs 9 agents with the Executor → Critic self-healing loop.", VIOLET),
    ("4  Model training", "FLAML 80/20 split, optimizing ROC-AUC or R² for the best model.", ACID),
    ("5  Explainability", "SHAP ranks features with the direction of impact + statistical tests.", MAGENTA),
    ("6  Outputs", "Model verdict, 16+ charts, a 3D graph and a board-ready report.", GOLD),
]
for i, (h, b, c) in enumerate(flow):
    x = 1.0 + (i % 3) * 3.83
    y = 2.4 + (i // 3) * 1.95
    box(s, x, y, 3.55, 1.7, PANEL, EDGE)
    box(s, x, y, 3.55, 0.1, c, None, radius=0.0)
    one(s, x + 0.3, y + 0.25, 3.0, 0.4, h, 14, c, True, HEAD)
    one(s, x + 0.3, y + 0.72, 3.05, 0.85, b, 11, MUTE, False, BODY)
one(s, 1.0, 6.75, 11.3, 0.4, "Every stage streams real-time status to the frontend via SSE — minutes, not days, from upload to insight.", 12, MUTE, False, BODY, CENTER)
footer(s, 14)

# ============================ 15 — LLM providers ============================
s = slide(INK)
hero_right(s, "art_brain.jpg", tag="HERO — AI brain", accent=VIOLET)
eyebrow(s, 1.0, 0.8, "AI models", VIOLET)
title(s, 1.0, 1.15, "Eight providers, one brain", w=6.9)
prov = [("Groq", "llama-3.3-70b", "FREE", ACID), ("OpenAI", "gpt-4o-mini", "token", MUTE),
        ("Anthropic", "claude (Sonnet)", "token", MUTE), ("DeepSeek", "deepseek-chat", "token", MUTE),
        ("Mistral", "mistral-small", "token", MUTE), ("Gemini", "gemini-2.0-flash", "FREE", ACID),
        ("OpenRouter", "deepseek-chat", "token", MUTE), ("Z.ai", "glm-4.6", "token", MUTE)]
box(s, 1.0, 2.1, 7.0, 3.5, PANEL, EDGE)
one(s, 1.25, 2.28, 2.4, 0.3, "PROVIDER", 10, MUTE, True, MONO)
one(s, 3.55, 2.28, 3.0, 0.3, "DEFAULT MODEL", 10, MUTE, True, MONO)
one(s, 6.55, 2.28, 1.2, 0.3, "COST", 10, MUTE, True, MONO)
for i, (p, m, cost, cc) in enumerate(prov):
    y = 2.66 + i * 0.355
    one(s, 1.25, y, 2.3, 0.32, p, 12, MIST, True, BODY)
    one(s, 3.55, y, 3.0, 0.32, m, 11, MUTE, False, MONO)
    one(s, 6.55, y, 1.2, 0.32, cost, 11, cc, cost == "FREE", MONO)
one(s, 1.0, 5.85, 7.0, 0.9, "A unified layer routes each task to the right model — capability matched to complexity. Priority: UI settings → per-role env → global env → MockLLM fallback (runs with no keys at all).", 11, MUTE, False, BODY)
footer(s, 15)

# ============================ 16 — Prompt engineering ============================
s = slide(INK2)
eyebrow(s, 1.0, 0.8, "AI models", GOLD)
title(s, 1.0, 1.15, "Prompt engineering & determinism")
one(s, 1.0, 2.0, 11.2, 0.8, "Six of the nine agents are LLM-powered. Each carries an engineered system prompt — an expert persona, a strict output schema and domain context — so results are structured, parseable and reproducible.", 13.5, MIST, False, BODY)
# code card
box(s, 1.0, 3.05, 5.6, 3.2, INK, EDGE)
one(s, 1.25, 3.22, 5.2, 0.3, "▌ SYSTEM PROMPT · Coder agent", 11, GOLD, True, MONO)
code = [('role: ', '"senior data scientist"', CYAN),
        ('task: ', "write Python for the plan", MIST),
        ('grounding: ', "RAG over pandas / sklearn", MIST),
        ('output: ', "FULL runnable script only", ACID),
        ('rule: ', "no prose, no markdown fences", CORAL),
        ('temp: ', "0.2  (deterministic)", GOLD)]
for i, (k, v, c) in enumerate(code):
    y = 3.7 + i * 0.42
    text(s, 1.3, y, 5.1, 0.4, [[(k, 12, MUTE, False, MONO), (v, 12, c, False, MONO)]])
pts = [("6 of 9 agents use LLMs", "Planner · Coder · Critic · Visualizer · Researcher · Reporter."),
       ("Temperature 0.2", "Low temperature across providers for deterministic output."),
       ("Strict JSON schema", "Machine-readable output → the UI renders charts directly."),
       ("MockLLM fallback", "Zero-config: the system runs with no API keys at all.")]
for i, (h, b) in enumerate(pts):
    y = 3.05 + i * 0.82
    box(s, 6.85, y, 5.48, 0.72, PANEL, EDGE)
    one(s, 7.08, y + 0.08, 5.1, 0.35, h, 12.5, GOLD, True, HEAD)
    one(s, 7.08, y + 0.4, 5.1, 0.3, b, 10, MUTE, False, BODY)
footer(s, 16)

# ============================ 17 — Nine task types ============================
s = slide(INK)
eyebrow(s, 1.0, 0.8, "Breadth", TEAL)
title(s, 1.0, 1.15, "Nine task types, auto-detected")
tasks = [
    ("Classification", "binary + multiclass", CYAN), ("Regression", "predict a number", AZURE),
    ("Clustering", "segments (k by silhouette)", IRIS), ("NLP + analytics", "topics · sentiment · keywords", MAGENTA),
    ("Anomaly detection", "Isolation Forest", CORAL), ("Dimensionality reduction", "PCA → 2-D map", VIOLET),
    ("Time-series", "Holt-Winters ETS + seasonality", GOLD), ("Survival analysis", "Kaplan-Meier + Cox PH", TEAL),
    ("Recommendation", "content-based similarity", ACID),
]
for i, (h, b, c) in enumerate(tasks):
    x = 1.0 + (i % 3) * 3.83
    y = 2.4 + (i // 3) * 1.42
    box(s, x, y, 3.55, 1.25, PANEL, EDGE)
    dot(s, x + 0.3, y + 0.3, 0.2, c)
    one(s, x + 0.66, y + 0.2, 2.8, 0.45, h, 14, WHITE, True, HEAD)
    one(s, x + 0.3, y + 0.68, 3.1, 0.45, b, 10.5, MUTE, False, MONO)
one(s, 1.0, 6.7, 11.3, 0.4, "Plus statistics, model-evaluation diagnostics, SMOTE, feature engineering, a what-if simulator and natural-language Q&A.", 11.5, MUTE, False, BODY, CENTER)
footer(s, 17)

# ============================ 18 — Statistics + evaluation ============================
s = slide(INK)
eyebrow(s, 1.0, 0.8, "Rigor", CYAN)
title(s, 1.0, 1.15, "Statistics + evaluation, not just a score")
box(s, 1.0, 2.25, 5.55, 4.2, PANEL, EDGE)
one(s, 1.3, 2.48, 5.0, 0.4, "Statistical significance", 15, CYAN, True, HEAD)
one(s, 1.3, 3.0, 5.0, 1.9, "Auto-picks the right test — Welch t / Mann-Whitney, ANOVA / Kruskal, chi-square, Pearson — each with a p-value, an effect size (Cohen's d, Cramér's V, |r|) and a plain-English verdict.", 12.5, MIST, False, BODY)
one(s, 1.3, 5.15, 5.0, 1.1, "So a weak dataset explains itself instead of looking broken — the honest answer, with the maths to back it.", 12.5, MUTE, False, BODY)
ev = ["Stratified k-fold cross-validation", "Confusion matrix + per-class P/R/F1",
      "ROC & precision-recall curves", "Calibration + learning curves",
      "Residuals (predicted vs actual)", "Model-comparison vs baselines"]
box(s, 6.78, 2.25, 5.55, 4.2, PANEL, EDGE)
one(s, 7.08, 2.48, 5.0, 0.4, "Model evaluation · diagnostics", 15, MAGENTA, True, HEAD)
for i, e in enumerate(ev):
    y = 3.08 + i * 0.53
    dot(s, 7.08, y + 0.08, 0.13, MAGENTA)
    one(s, 7.34, y, 4.9, 0.45, e, 12.5, MIST, False, BODY)
footer(s, 18)

# ============================ 19 — Security ============================
s = slide(INK)
hero_right(s, "art_security.jpg", tag="HERO — shield", accent=CYAN)
eyebrow(s, 1.0, 0.8, "Security", CYAN)
title(s, 1.0, 1.15, "Sandboxing & secure execution", w=6.9)
one(s, 1.0, 2.0, 7.0, 0.85, "AI-generated code is never trusted blindly. Helix runs it behind a defense-in-depth dual sandbox and auto-selects the strongest layer available.", 13, MIST, False, BODY)
sec = [("RestrictedPython", "Default layer. AST-level restriction blocks dangerous imports — os, sys, socket — and unsafe attribute access before code runs.", AZURE),
       ("E2B microVM", "Production layer. Full VM isolation: no host filesystem or network, with a hard execution timeout.", VIOLET),
       ("Input sanitization", "CSV files are scanned and column names / values sanitized before any exposure to an LLM.", CORAL)]
for i, (h, b, c) in enumerate(sec):
    y = 3.05 + i * 1.18
    box(s, 1.0, y, 7.0, 1.05, PANEL, EDGE)
    box(s, 1.0, y, 0.12, 1.05, c, None, radius=0.0)
    one(s, 1.35, y + 0.13, 6.5, 0.4, h, 14, c, True, HEAD)
    one(s, 1.35, y + 0.52, 6.5, 0.45, b, 11, MUTE, False, BODY)
footer(s, 19)

# ============================ 20 — Resilience + SSE ============================
s = slide(INK2)
eyebrow(s, 1.0, 0.8, "Resilience", ACID)
title(s, 1.0, 1.15, "Graceful degradation — it never breaks")
one(s, 1.0, 2.0, 11.2, 0.5, "Every external dependency has a fallback. If something is missing or fails, Helix degrades one level instead of breaking — the app always returns a complete response.", 13, MIST, False, BODY)
levels = [("LLM", "MockLLM", VIOLET), ("RAG", "keyword match", CYAN), ("Tavily", "DuckDuckGo", AZURE),
          ("E2B", "RestrictedPython", CORAL), ("Pipeline error", "mock events", GOLD), ("Backend down", "client simulation", ACID)]
for i, (a, b, c) in enumerate(levels):
    x = 1.0 + (i % 3) * 3.83
    y = 2.85 + (i // 3) * 1.35
    box(s, x, y, 3.55, 1.15, PANEL, EDGE)
    one(s, x + 0.25, y + 0.16, 3.1, 0.35, f"LEVEL {i + 1}", 10, MUTE, True, MONO)
    text(s, x + 0.25, y + 0.52, 3.1, 0.5, [[(a + "  ", 13, MIST, True, BODY), ("→ " + b, 12, c, False, MONO)]])
one(s, 1.0, 6.6, 11.3, 0.4, "Real-time SSE: an asyncio.Queue bus + anti-buffering headers push every agent event to the UI with sub-second latency.", 12, MUTE, False, BODY, CENTER)
footer(s, 20)

# ============================ 21 — Deployment + comparison ============================
s = slide(INK)
eyebrow(s, 1.0, 0.8, "Deployment & edge", IRIS)
title(s, 1.0, 1.15, "Deploy anywhere — how Helix compares")
deploys = [("Vercel + Hugging Face", "Production — global CDN frontend + containerized backend.", CYAN),
           ("Docker Compose", "Local — one command: docker compose up --build.", AZURE),
           ("Google Colab", "Quick demo — pre-built notebooks, zero install.", ACID)]
for i, (h, b, c) in enumerate(deploys):
    y = 2.25 + i * 1.4
    box(s, 1.0, y, 5.4, 1.25, PANEL, EDGE)
    dot(s, 1.3, y + 0.4, 0.2, c)
    one(s, 1.7, y + 0.18, 4.6, 0.4, h, 14, WHITE, True, HEAD)
    one(s, 1.7, y + 0.62, 4.6, 0.5, b, 11, MUTE, False, BODY)
box(s, 6.75, 2.25, 5.58, 4.4, PANEL, EDGE)
one(s, 7.0, 2.42, 3.0, 0.3, "CAPABILITY", 10.5, MUTE, True, MONO)
one(s, 10.25, 2.42, 1.1, 0.3, "HELIX", 10.5, ACID, True, MONO, CENTER)
one(s, 11.4, 2.42, 1.1, 0.3, "OTHERS", 10.5, MUTE, True, MONO, CENTER)
caps = [("Plain-English input", "✕"), ("Self-healing code (5×)", "✕"), ("9 collaborating agents", "✕"),
        ("Real ML training (FLAML)", "partial"), ("SHAP explainability", "manual"),
        ("Live web research", "✕"), ("Business report", "✕"), ("9 task types", "partial"),
        ("3D knowledge graph", "✕")]
for i, (cap, other) in enumerate(caps):
    y = 2.86 + i * 0.41
    one(s, 7.0, y, 3.2, 0.35, cap, 11.5, MIST, False, BODY, anchor=MID)
    one(s, 10.25, y, 1.1, 0.35, "✓", 13, ACID, True, BODY, CENTER, MID)
    one(s, 11.4, y, 1.1, 0.35, other, 11, CORAL if other == "✕" else GOLD, False, BODY, CENTER, MID)
footer(s, 21)

# ============================ 22 — Conclusion ============================
s = slide(INK)
hero_right(s, "art_conclusion.jpg", tag="NANOBANANA — closing", accent=CYAN)
eyebrow(s, 1.0, 0.85, "Conclusion", CYAN)
title(s, 1.0, 1.25, "Data science, reimagined", w=6.9)
one(s, 1.0, 2.15, 7.0, 1.2, "Helix proves a nine-agent pipeline can run the entire data-science workflow autonomously — from a plain-English goal to a board-ready report with real ML, SHAP explainability, 16+ charts and web-grounded insight.", 13.5, MIST, False, BODY)
pillars = [("Hybrid LLM hosting", "8 providers, task-matched, MockLLM fallback", VIOLET),
           ("Dual-sandbox execution", "RestrictedPython + E2B, auto-selected", AZURE),
           ("LangGraph orchestration", "shared state, conditional self-healing", ACID),
           ("Stateless architecture", "no DB, no sessions — failure classes gone", GOLD)]
for i, (h, b, c) in enumerate(pillars):
    y = 3.55 + i * 0.62
    dot(s, 1.0, y + 0.1, 0.16, c)
    text(s, 1.3, y, 6.7, 0.5, [[(h + "  ", 12.5, WHITE, True, HEAD), (b, 11.5, MUTE, False, BODY)]])
text(s, 1.0, 6.15, 7.0, 0.5, [[("Days  →  Minutes.", 22, CYAN, True, HEAD)]])
one(s, 1.0, 6.8, 7.2, 0.35, "helix-henna.vercel.app   ·   github.com/Srujan29112001/helix   ·   IIIT-H · June 2026", 10.5, MUTE, False, MONO)
footer(s, 22)

# ----------------------------- save -----------------------------
out_path = os.path.join(HERE, "Helix_Project_Deck.pptx")
try:
    prs.save(out_path)
except PermissionError:
    out_path = os.path.join(HERE, "Helix_Project_Deck_new.pptx")
    prs.save(out_path)
    print("NOTE: primary file was locked (open in PowerPoint?) — wrote", os.path.basename(out_path))
n = len(prs.slides._sldIdLst)
imgs = ["art_hero.jpg", "05_landing.png", "art_problem.jpg", "diagram_architecture.png",
        "06_studio.png", "10_pipeline.png", "11_results_telco.png", "art_backend.jpg",
        "art_brain.jpg", "art_security.jpg", "art_conclusion.jpg"]
have = sum(os.path.exists(os.path.join(ASSETS, f)) for f in imgs)
print(f"saved {out_path} · {n} slides · {have}/{len(imgs)} images embedded")
