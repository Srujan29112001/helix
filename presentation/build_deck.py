"""Builds the Helix capstone presentation deck (dark, brand-matched).
Run:  python presentation/build_deck.py  ->  presentation/Helix_Capstone.pptx
"""
import os

from pptx import Presentation
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
from pptx.enum.text import MSO_ANCHOR, PP_ALIGN
from pptx.util import Inches, Pt

# ---- palette ----
INK = RGBColor(0x04, 0x06, 0x0F)
INK2 = RGBColor(0x09, 0x0E, 0x1A)
PANEL = RGBColor(0x0C, 0x12, 0x24)
EDGE = RGBColor(0x1A, 0x24, 0x40)
CYAN = RGBColor(0x25, 0xD7, 0xF0)
VIOLET = RGBColor(0x8B, 0x5C, 0xF6)
AZURE = RGBColor(0x38, 0xBD, 0xF8)
IRIS = RGBColor(0xA7, 0x8B, 0xFA)
MAGENTA = RGBColor(0xE8, 0x79, 0xF9)
ACID = RGBColor(0x9A, 0xE6, 0x4A)
GOLD = RGBColor(0xFB, 0xBF, 0x24)
CORAL = RGBColor(0xFB, 0x71, 0x85)
WHITE = RGBColor(0xFF, 0xFF, 0xFF)
MIST = RGBColor(0xC9, 0xD6, 0xEF)
MUTE = RGBColor(0x76, 0x85, 0x9F)

HEAD = "Segoe UI Semibold"
BODY = "Segoe UI"
MONO = "Consolas"

prs = Presentation()
prs.slide_width = Inches(13.333)
prs.slide_height = Inches(7.5)
BLANK = prs.slide_layouts[6]
SW, SH = 13.333, 7.5


def slide(bg=INK):
    s = prs.slides.add_slide(BLANK)
    r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
    r.fill.solid()
    r.fill.fore_color.rgb = bg
    r.line.fill.background()
    r.shadow.inherit = False
    return s


def text(s, x, y, w, h, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP, space=2):
    tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(h))
    tf = tb.text_frame
    tf.word_wrap = True
    tf.vertical_anchor = anchor
    tf.margin_left = tf.margin_right = tf.margin_top = tf.margin_bottom = 0
    if isinstance(runs[0], tuple):
        runs = [runs]
    for i, para in enumerate(runs):
        p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
        p.alignment = align
        p.space_after = Pt(space)
        for t, size, color, bold, font in para:
            r = p.add_run()
            r.text = t
            r.font.size = Pt(size)
            r.font.color.rgb = color
            r.font.bold = bold
            r.font.name = font
    return tb


def one(s, x, y, w, h, t, size, color, bold=False, font=BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
    return text(s, x, y, w, h, [[(t, size, color, bold, font)]], align=align, anchor=anchor)


def box(s, x, y, w, h, fill=PANEL, line=EDGE, line_w=1.0, radius=0.08):
    shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(h))
    shp.adjustments[0] = radius
    shp.fill.solid()
    shp.fill.fore_color.rgb = fill
    if line is None:
        shp.line.fill.background()
    else:
        shp.line.color.rgb = line
        shp.line.width = Pt(line_w)
    shp.shadow.inherit = False
    return shp


def dot(s, x, y, d, color):
    c = s.shapes.add_shape(MSO_SHAPE.OVAL, Inches(x), Inches(y), Inches(d), Inches(d))
    c.fill.solid()
    c.fill.fore_color.rgb = color
    c.line.fill.background()
    c.shadow.inherit = False
    return c


def eyebrow(s, x, y, label, color=CYAN):
    dot(s, x, y + 0.05, 0.12, color)
    one(s, x + 0.22, y - 0.05, 6, 0.4, label.upper(), 12, MUTE, True, MONO)


def title(s, x, y, t, w=11.5):
    one(s, x, y, w, 1.1, t, 34, WHITE, True, HEAD)


AGENTS = [
    ("Planner", VIOLET), ("Coder", CYAN), ("Executor", AZURE), ("Critic", GOLD),
    ("AutoML", ACID), ("Explainer", MAGENTA), ("Reporter", CORAL),
]

# ============================ Slide 1 — Title ============================
s = slide(INK)
# motif: agent-color dots
for i, (_, c) in enumerate(AGENTS):
    dot(s, 1.0 + i * 0.42, 1.7, 0.22, c)
text(s, 1.0, 2.5, 11.5, 2.2, [
    [("Helix", 76, WHITE, True, HEAD), (".", 76, CYAN, True, HEAD)],
])
text(s, 1.02, 3.85, 11.5, 1.0, [[("Autonomous Data Science Agent ", 26, MIST, False, BODY),
                                  ("with Self-Correcting Code Execution", 26, CYAN, True, BODY)]])
one(s, 1.05, 5.0, 11, 0.5, "From a CSV and a sentence to a business-ready report — automatically.", 16, MUTE, False, BODY)
box(s, 1.0, 6.2, 4.2, 0.7, INK2, EDGE)
one(s, 1.25, 6.32, 4, 0.45, "Capstone Project  ·  IIIT-H", 14, MIST, False, MONO, anchor=MSO_ANCHOR.MIDDLE)

# ============================ Slide 2 — Problem ============================
s = slide(INK)
eyebrow(s, 1.0, 0.85, "The problem", CORAL)
title(s, 1.0, 1.2, "Insight is trapped behind a bottleneck")
probs = [
    ("Scarce expertise", "Business teams queue for a data scientist to answer even simple questions — days of delay.", CORAL),
    ("Repetitive grind", "60–70% of analyst time goes to cleaning CSVs, handling missing values and boilerplate EDA.", GOLD),
    ("Brittle iteration", "Code breaks, you read the traceback, patch, rerun, repeat — painful for non-experts.", IRIS),
]
for i, (h, b, c) in enumerate(probs):
    x = 1.0 + i * 3.85
    box(s, x, 2.6, 3.55, 3.3, PANEL, EDGE)
    dot(s, x + 0.35, 2.95, 0.28, c)
    one(s, x + 0.35, 3.5, 3, 0.5, h, 19, WHITE, True, HEAD)
    one(s, x + 0.35, 4.1, 2.95, 1.6, b, 13.5, MUTE, False, BODY)

# ============================ Slide 3 — Solution ============================
s = slide(INK)
eyebrow(s, 1.0, 0.85, "The solution")
title(s, 1.0, 1.2, "An AI team that does the whole workflow")
flow = ["Upload CSV", "Plain-English goal", "7 agents work", "Self-correct", "Report"]
for i, step in enumerate(flow):
    x = 1.0 + i * 2.4
    box(s, x, 2.9, 2.1, 0.95, INK2, CYAN if i in (0, 4) else EDGE)
    one(s, x, 2.9, 2.1, 0.95, step, 14, MIST, True, BODY, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
    if i < len(flow) - 1:
        one(s, x + 2.05, 2.9, 0.4, 0.95, "›", 26, MUTE, True, BODY, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
box(s, 1.0, 4.5, 11.3, 1.7, PANEL, EDGE)
text(s, 1.4, 4.85, 10.6, 1.1, [[('“Predict which customers will churn.”', 18, WHITE, True, BODY)],
                               [("→  82–84% ROC-AUC, top drivers = contract type & tenure, with a plain-English report — in minutes.", 15, ACID, False, BODY)]], space=8)

# ====================== Slide 4 — Architecture ======================
s = slide(INK)
eyebrow(s, 1.0, 0.7, "Architecture", VIOLET)
title(s, 1.0, 1.05, "Seven agents, orchestrated by LangGraph")
for i, (name, c) in enumerate(AGENTS):
    x = 0.75 + i * 1.78
    box(s, x, 2.9, 1.5, 1.0, INK2, c, 1.5)
    one(s, x, 3.05, 1.5, 0.4, name, 12.5, WHITE, True, BODY, PP_ALIGN.CENTER)
    one(s, x, 3.45, 1.5, 0.35, "agent", 9, MUTE, False, MONO, PP_ALIGN.CENTER)
    if i < len(AGENTS) - 1:
        one(s, x + 1.5, 2.9, 0.28, 1.0, "›", 20, MUTE, True, BODY, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
box(s, 3.2, 4.3, 4.7, 0.6, INK2, GOLD)
one(s, 3.2, 4.3, 4.7, 0.6, "↺  self-correction loop · Executor ⇄ Critic · ≤ 5 retries", 12.5, GOLD, True, MONO, PP_ALIGN.CENTER, MSO_ANCHOR.MIDDLE)
roles = "Planner (CoT plan) · Coder (RAG-grounded) · Executor (sandbox) · Critic (self-heal) · AutoML (FLAML) · Explainer (SHAP) · Reporter (LLM)"
one(s, 1.0, 5.6, 11.3, 0.8, roles, 13, MUTE, False, BODY, PP_ALIGN.CENTER)

# ====================== Slide 5 — Self-correcting execution ======================
s = slide(INK2)
eyebrow(s, 1.0, 0.85, "The headline capability", GOLD)
title(s, 1.0, 1.2, "Self-correcting code execution")
steps = [
    ("Coder writes Python", "df.mean(numeric_only=False)", CYAN),
    ("Executor runs it (sandbox)", "TypeError: could not convert…", CORAL),
    ("Critic reads the traceback", "fix → numeric_only=True", GOLD),
    ("Re-run → clean output", "tenure 32.37 · charges 64.76", ACID),
]
for i, (h, code, c) in enumerate(steps):
    y = 2.6 + i * 0.95
    dot(s, 1.0, y + 0.12, 0.3, c)
    one(s, 1.5, y, 5, 0.5, h, 16, WHITE, True, BODY)
    box(s, 6.7, y - 0.05, 5.6, 0.6, INK, EDGE)
    one(s, 6.9, y - 0.05, 5.2, 0.6, code, 13, c, False, MONO, anchor=MSO_ANCHOR.MIDDLE)
one(s, 1.0, 6.6, 11, 0.5, "RestrictedPython sandbox — file-system, network and dangerous imports all blocked.", 13, MUTE, False, BODY)

# ====================== Slide 6 — Capabilities ======================
s = slide(INK)
eyebrow(s, 1.0, 0.85, "Capabilities", ACID)
title(s, 1.0, 1.2, "Real ML, sandboxed and explainable")
caps = [
    ("FLAML AutoML", "Auto-selects & tunes the best model under a time budget", ACID),
    ("SHAP explainability", "Quantifies what drives every prediction", MAGENTA),
    ("RAG grounding", "ChromaDB + embeddings ground the Coder's code", CYAN),
    ("Secure sandbox", "RestrictedPython blocks fs / network / imports", AZURE),
    ("4 task types", "Classification · regression · clustering · NLP", IRIS),
    ("Any CSV", "Auto-cleans messy data and detects the task", GOLD),
]
for i, (h, b, c) in enumerate(caps):
    x = 1.0 + (i % 3) * 3.85
    y = 2.5 + (i // 3) * 2.0
    box(s, x, y, 3.55, 1.75, PANEL, EDGE)
    dot(s, x + 0.32, y + 0.32, 0.22, c)
    one(s, x + 0.72, y + 0.26, 2.7, 0.5, h, 16, WHITE, True, HEAD)
    one(s, x + 0.32, y + 0.85, 3.0, 0.8, b, 12.5, MUTE, False, BODY)

# ====================== Slide 7 — Results / Evaluation ======================
s = slide(INK2)
eyebrow(s, 1.0, 0.85, "Evaluation", CYAN)
title(s, 1.0, 1.2, "Measured across all four task types")
stats = [("100%", "task detection", ACID), ("100%", "self-correction", CYAN), ("100%", "sandbox security", GOLD), ("7.1s", "avg runtime", IRIS)]
for i, (v, l, c) in enumerate(stats):
    x = 1.0 + i * 2.85
    box(s, x, 2.5, 2.6, 1.5, PANEL, EDGE)
    one(s, x, 2.6, 2.6, 0.8, v, 38, c, True, HEAD, PP_ALIGN.CENTER)
    one(s, x, 3.5, 2.6, 0.4, l, 12, MUTE, False, MONO, PP_ALIGN.CENTER)
rows = [
    ("Telco Churn", "Binary classification", "Random Forest", "0.84 ROC-AUC"),
    ("Sales (synthetic)", "Regression", "Extra Trees", "0.95 R²"),
    ("Segments", "Clustering", "KMeans (k=4)", "4 segments"),
    ("Reviews", "NLP + analytics", "LightGBM", "100% accuracy"),
]
box(s, 1.0, 4.35, 11.3, 2.5, PANEL, EDGE)
hdr = ["Dataset", "Detected task", "Best model", "Result"]
for j, htxt in enumerate(hdr):
    one(s, 1.4 + j * 2.8, 4.55, 2.7, 0.4, htxt, 12, MUTE, True, MONO)
for i, row in enumerate(rows):
    y = 5.05 + i * 0.43
    for j, cell in enumerate(row):
        col = ACID if j == 3 else MIST
        one(s, 1.4 + j * 2.8, y, 2.7, 0.4, cell, 12.5, col, j == 3, BODY)

# ====================== Slide 8 — Tech stack ======================
s = slide(INK)
eyebrow(s, 1.0, 0.85, "Tech stack", IRIS)
title(s, 1.0, 1.2, "Open models, production tooling")
groups = [
    ("Agents / LLMs", "LangGraph · LangChain · DeepSeek · Mistral · Groq", VIOLET),
    ("Execution", "RestrictedPython sandbox · FastAPI (SSE)", AZURE),
    ("ML / AutoML", "scikit-learn · FLAML · pandas · NumPy", ACID),
    ("Explainability", "SHAP", MAGENTA),
    ("Retrieval (RAG)", "ChromaDB · sentence embeddings (MiniLM)", CYAN),
    ("Interface", "Next.js · React · Tailwind · Gradio (Colab)", GOLD),
]
for i, (h, b, c) in enumerate(groups):
    x = 1.0 + (i % 2) * 5.85
    y = 2.4 + (i // 2) * 1.45
    box(s, x, y, 5.5, 1.2, PANEL, EDGE)
    dot(s, x + 0.32, y + 0.45, 0.2, c)
    one(s, x + 0.7, y + 0.18, 4.6, 0.45, h, 15, WHITE, True, HEAD)
    one(s, x + 0.7, y + 0.62, 4.6, 0.5, b, 12, MUTE, False, MONO)

# ====================== Slide 9 — Outcomes ======================
s = slide(INK)
eyebrow(s, 1.0, 0.85, "Impact", CORAL)
title(s, 1.0, 1.2, "From days to minutes")
out = [
    ("Faster insight", "Analyses that took days now take minutes.", CYAN),
    ("Less dependency", "Non-technical users run analyses themselves.", ACID),
    ("Explainable", "Every result ships with SHAP drivers + a narrative.", MAGENTA),
    ("Trustworthy", "Sandboxed, self-correcting, and reproducible.", GOLD),
]
for i, (h, b, c) in enumerate(out):
    x = 1.0 + (i % 2) * 5.85
    y = 2.5 + (i // 2) * 1.9
    box(s, x, y, 5.5, 1.6, PANEL, EDGE)
    one(s, x + 0.4, y + 0.28, 4.7, 0.5, h, 19, c, True, HEAD)
    one(s, x + 0.4, y + 0.85, 4.7, 0.6, b, 13.5, MUTE, False, BODY)

# ====================== Slide 10 — Closing ======================
s = slide(INK)
for i, (_, c) in enumerate(AGENTS):
    dot(s, 1.0 + i * 0.42, 1.8, 0.2, c)
text(s, 1.0, 2.6, 11, 1.4, [[("From a CSV to a decision —", 40, WHITE, True, HEAD)],
                            [("autonomously.", 40, CYAN, True, HEAD)]], space=4)
one(s, 1.02, 4.5, 11, 0.5, "Live demo:  localhost:3000/studio   ·   Colab:  Helix_RealData_Colab.ipynb", 15, MIST, False, MONO)
one(s, 1.02, 5.2, 11, 0.5, "github · capstone · IIIT-H", 13, MUTE, False, MONO)
one(s, 1.0, 6.4, 11, 0.5, "Thank you.", 18, MUTE, True, BODY)

out_path = os.path.join(os.path.dirname(__file__), "Helix_Capstone.pptx")
prs.save(out_path)
print("saved", out_path, "·", len(prs.slides._sldIdLst), "slides")
