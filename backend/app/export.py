"""Export an analysis result to a polished PowerPoint deck (.pptx) or Markdown."""

from __future__ import annotations

import io
from typing import Any


def _g(results: dict, key: str, default=None):
    return results.get(key, default) if isinstance(results, dict) else default


def build_markdown(results: dict, goal: str = "", dataset: str = "dataset") -> bytes:
    h = _g(results, "headline", {}) or {}
    lines = [
        f"# Helix Analysis — {dataset}",
        "",
        f"**Goal:** {goal or '—'}",
        f"**Task:** {_g(results, 'taskLabel', '?')}  ·  **Best model:** {_g(results, 'bestModel', '?')}",
        f"**Headline:** {h.get('value', '')} {h.get('label', '')}",
        "",
        "## Model performance",
        "",
        "| Metric | Value |",
        "|--------|-------|",
    ]
    for m in _g(results, "metrics", []) or []:
        lines.append(f"| {m.get('label')} | {m.get('value')} |")
    v = _g(results, "_verdict") or {}
    if v:
        lines += ["", f"**Model quality:** {v.get('label')} — {v.get('detail')}"]
    lines += ["", "## Key drivers", ""]
    for b in _g(results, "bars", []) or []:
        sign = b.get("sign", 0)
        eff = " (raises)" if sign and sign > 0 else " (lowers)" if sign and sign < 0 else ""
        lines.append(f"- **{b.get('label')}** — importance {round(float(b.get('value', 0)), 2)}{eff}")
    st = _g(results, "_stats_tests") or []
    if st:
        lines += ["", "## Statistical significance", ""]
        for t in st[:8]:
            p = "<0.001" if t.get("p", 1) < 0.001 else f"{t.get('p', 0):.3f}"
            lines.append(f"- {t.get('feature')}: {t.get('test')}, p={p} — {'significant' if t.get('significant') else 'n.s.'}")
    lines += ["", "## Business report", ""]
    for p in _g(results, "report", []) or []:
        lines += [str(p), ""]
    rec = _g(results, "recommendation")
    if rec:
        lines += ["## Recommendation", "", f"> {rec}", ""]
    return ("\n".join(lines)).encode("utf-8")


def build_pptx(results: dict, goal: str = "", dataset: str = "dataset") -> bytes:
    from pptx import Presentation
    from pptx.dml.color import RGBColor
    from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
    from pptx.enum.shapes import MSO_SHAPE
    from pptx.util import Inches, Pt

    INK = RGBColor(0x06, 0x09, 0x12)
    PANEL = RGBColor(0x0C, 0x12, 0x24)
    EDGE = RGBColor(0x1A, 0x24, 0x40)
    CYAN = RGBColor(0x25, 0xD7, 0xF0)
    ACID = RGBColor(0x9A, 0xE6, 0x4A)
    WHITE = RGBColor(0xFF, 0xFF, 0xFF)
    MIST = RGBColor(0xC9, 0xD6, 0xEF)
    MUTE = RGBColor(0x76, 0x85, 0x9F)
    HEAD, BODY, MONO = "Segoe UI Semibold", "Segoe UI", "Consolas"

    prs = Presentation()
    prs.slide_width, prs.slide_height = Inches(13.333), Inches(7.5)
    BLANK = prs.slide_layouts[6]

    def slide():
        s = prs.slides.add_slide(BLANK)
        r = s.shapes.add_shape(MSO_SHAPE.RECTANGLE, 0, 0, prs.slide_width, prs.slide_height)
        r.fill.solid(); r.fill.fore_color.rgb = INK; r.line.fill.background(); r.shadow.inherit = False
        return s

    def tx(s, x, y, w, hgt, runs, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        tb = s.shapes.add_textbox(Inches(x), Inches(y), Inches(w), Inches(hgt))
        tf = tb.text_frame; tf.word_wrap = True; tf.vertical_anchor = anchor
        for i, para in enumerate(runs):
            p = tf.paragraphs[0] if i == 0 else tf.add_paragraph()
            p.alignment = align; p.space_after = Pt(4)
            for t, sz, col, bold, font in para:
                run = p.add_run(); run.text = t; run.font.size = Pt(sz)
                run.font.color.rgb = col; run.font.bold = bold; run.font.name = font
        return tb

    def one(s, x, y, w, hgt, t, sz, col, bold=False, font=BODY, align=PP_ALIGN.LEFT, anchor=MSO_ANCHOR.TOP):
        return tx(s, x, y, w, hgt, [[(t, sz, col, bold, font)]], align, anchor)

    def box(s, x, y, w, hgt, line=EDGE):
        shp = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(x), Inches(y), Inches(w), Inches(hgt))
        shp.adjustments[0] = 0.06; shp.fill.solid(); shp.fill.fore_color.rgb = PANEL
        shp.line.color.rgb = line; shp.line.width = Pt(1); shp.shadow.inherit = False
        return shp

    h = _g(results, "headline", {}) or {}

    # 1 — title
    s = slide()
    tx(s, 1, 2.4, 11, 1.4, [[("Helix", 60, WHITE, True, HEAD), (" Analysis", 60, CYAN, True, HEAD)]])
    one(s, 1.03, 3.7, 11, 0.6, str(dataset), 22, MIST, False, BODY)
    one(s, 1.03, 4.4, 11, 0.6, (goal or "")[:120], 15, MUTE, False, BODY)
    one(s, 1.03, 5.6, 11, 0.5,
        f"{_g(results, 'taskLabel', '')}  ·  best model {_g(results, 'bestModel', '')}  ·  {h.get('value','')} {h.get('label','')}",
        14, ACID, True, MONO)

    # 2 — performance
    s = slide()
    one(s, 1, 0.7, 11, 0.8, "Model performance", 30, WHITE, True, HEAD)
    mets = _g(results, "metrics", []) or []
    for i, m in enumerate(mets[:5]):
        x = 1 + i * 2.4
        box(s, x, 2.0, 2.2, 1.5)
        one(s, x, 2.15, 2.2, 0.7, str(m.get("value", "")), 30, CYAN, True, HEAD, PP_ALIGN.CENTER)
        one(s, x, 3.05, 2.2, 0.4, str(m.get("label", "")), 12, MUTE, False, MONO, PP_ALIGN.CENTER)
    v = _g(results, "_verdict") or {}
    if v:
        box(s, 1, 4.0, 11.3, 1.6)
        one(s, 1.3, 4.2, 10.7, 0.5, f"Model quality · {v.get('label','')}", 16, ACID, True, HEAD)
        one(s, 1.3, 4.8, 10.7, 0.8, str(v.get("detail", ""))[:200], 13, MIST, False, BODY)

    # 3 — drivers
    s = slide()
    one(s, 1, 0.7, 11, 0.8, "Key drivers", 30, WHITE, True, HEAD)
    bars = _g(results, "bars", []) or []
    for i, b in enumerate(bars[:7]):
        y = 1.9 + i * 0.7
        sign = b.get("sign", 0)
        eff = "↑ raises" if sign and sign > 0 else "↓ lowers" if sign and sign < 0 else ""
        one(s, 1.2, y, 5, 0.5, str(b.get("label", "")), 16, MIST, False, BODY)
        bw = max(0.2, float(b.get("value", 0))) * 4.5
        bar = s.shapes.add_shape(MSO_SHAPE.ROUNDED_RECTANGLE, Inches(6.4), Inches(y + 0.05), Inches(bw), Inches(0.32))
        bar.fill.solid(); bar.fill.fore_color.rgb = CYAN; bar.line.fill.background(); bar.shadow.inherit = False
        one(s, 11.2, y, 1.6, 0.5, eff, 12, MUTE, False, MONO)

    # 4 — report + recommendation
    s = slide()
    one(s, 1, 0.6, 11, 0.7, "Findings & recommendation", 28, WHITE, True, HEAD)
    paras = [[(str(p), 13, MIST, False, BODY)] for p in (_g(results, "report", []) or [])[:5]]
    if paras:
        tx(s, 1, 1.5, 11.3, 4.0, paras)
    rec = _g(results, "recommendation")
    if rec:
        box(s, 1, 5.7, 11.3, 1.3, CYAN)
        one(s, 1.3, 5.85, 10.7, 0.4, "RECOMMENDATION", 11, MUTE, True, MONO)
        one(s, 1.3, 6.25, 10.7, 0.7, str(rec)[:240], 15, WHITE, True, BODY)

    buf = io.BytesIO()
    prs.save(buf)
    return buf.getvalue()
